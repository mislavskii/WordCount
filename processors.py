from pptx import Presentation
from docx import Document
# from googletrans import Translator

class Processor:
    RU_CHARS = 'абвгдеёжзийклмнопрстуфхцчшщыэюя'
                            
    def get_russian_words(self):
        self.russian_words = []
        self.other_words = []
        for paragraph in self.paragraphs: # type: ignore
            words = paragraph.lower().replace('\xa0', ' ').split()
            for word in words:
                word = word.strip()
                if word:
                    self.russian_words.append(word
                        ) if any(char in self.RU_CHARS for char in word
                        ) else self.other_words.append(word)
            
    def translate(self):
        for paragraph in self.paragraphs: # type: ignore
            print(paragraph)
            translation = input()

    def find(self, string):
        matches = [p for p in self.paragraphs if string.lower() in p.lower()]
        return matches
    

class PPTProcessor(Processor):

    def __init__(self, path):
        self.document = Presentation(path)
        self.get_paragraphs()
        
    def get_paragraphs(self):
        self.paragraphs = []
        for slide in self.document.slides:
            for shape in slide.shapes:
        # scanning text frames        
                if shape.has_text_frame: 
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            self.paragraphs.append(text)
        # scanning tables                    
                if shape.has_table: 
                    for cell in shape.table.iter_cells():
                        for paragraph in cell.text_frame.paragraphs:
                            self.paragraphs.append(paragraph.text)

            
class DocProcessor(Processor):
    
    def __init__(self, path):
        self.document = Document(path)
        self.get_paragraphs()

    def get_paragraphs(self):
        self.paragraphs = []
        for paragraph in self.document.paragraphs:
            text = paragraph.text.strip()
            if text:
                self.paragraphs.append(text)
        for table in self.document.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        text = paragraph.text.strip()
                        if text and text not in self.paragraphs:
                            self.paragraphs.append(text)