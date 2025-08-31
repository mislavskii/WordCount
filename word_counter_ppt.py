# import os
import string as s
from pptx import Presentation as pres


def count(file):
    """counts Russian words in a ppt file"""
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    prs = pres(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # scanning text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            # scanning tables
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
    # Counting words
    tot_word_count = 0
    rus_word_count = 0
    eng = []
    ru = []
    ru_chars = 'абвгдеёжзийклмнопрстуфхцчшщыэюя'

    for string in text_runs:
        words = string.strip().split()
        for word in words:
            check = False
            tot_word_count += 1
            for char in ru_chars:
                if char in word.lower():
                    rus_word_count += 1
                    ru.append(word)
                    check = True
                    break
                else:
                    continue
            if not check:
                eng.append(word)

    return len(ru)


class WordCounter:
    """counts Russian and English words in a PPT file"""
    ru_chars = 'абвгдеёжзийклмнопрстуфхцчшщыэюя'

    def __init__(self, path):
        self.presentation = pres(path)
        self.text_runs = []
        self.ru = []
        self.en = []
        self.dump = []
        self.ru_count = None
        self.en_count = None

    def get_texts(self):
        # scanning the presentation for text runs and collecting them into a list of strings
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                # scanning text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            self.text_runs.append(run.text)
                # scanning tables
                if shape.has_table:
                    for cell in shape.table.iter_cells():
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                self.text_runs.append(run.text)

    def count_words(self):
        self.get_texts()
        for string in self.text_runs:
            words = string.strip().split()
            for word in words:
                for char in word.lower():
                    if char in self.ru_chars:
                        self.ru.append(word)
                        break
                    elif char in s.ascii_lowercase:
                        self.en.append(word)
                        break
                    else:
                        continue
        self.en_count = len(self.en)  # number of words with roman chars and no cyrillic chars
        self.ru_count = len(self.ru)  # number of words with cyrillic chars and possibly some romans


class ListCounter:  # TODO: rewrite based on WordCounter
    """counts russian and english words in multiple files provided as a list of path-like objects"""
    word_count = 0

    def __init__(self, input_list):  # takes a list of file paths
        self.input_list = input_list

    def get_output_list(self):
        for file in self.input_list:
            pass


class PairedCounter:
    rate = 1.3
    amount = None
    lang = None

    def __init__(self, file1, file2):  # takes two ppt files as path-like objects
        self.file1 = file1
        self.file2 = file2
        self.lang = input('Language (en or ru): ')
        self.count1 = WordCounter(self.file1)
        self.count2 = WordCounter(self.file2)
        self.count1.count_words()
        self.count2.count_words()

    def calculate(self):  # calculates the job cost based on english word counts
        if self.lang == 'ru':
            self.amount = self.rate * abs(self.count1.ru_count - self.count2.ru_count)
        elif self.lang == 'en':
            self.amount = self.rate * abs(self.count1.en_count - self.count2.en_count)
        else:
            print('Wrong language. Go to hell!')
        return self.amount

